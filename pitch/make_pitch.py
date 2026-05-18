"""
Generates two deliverables from one source of truth:
  - <BRAND>_Business_Plan.docx
  - <BRAND>_Pitch_Deck.pptx

Concept v4 (2026-05-18):
  English-only edition. All prose, tables, and slide content in English.
  Validation-first thesis. "Second brain that validates."
  - 4-tier validation system (L1 Self / L2 Peer / L3 System / L4 Expert)
  - Obsidian as primary competitor, foregrounded in deck
  - Wikipedia / Community Notes as institutional validation precedents
  - Staged funding model (Phase 0 bootstrap -> Phase 1 Korean seed -> Phase 2)

Run:
    .venv/Scripts/python.exe make_pitch.py
"""
from pathlib import Path
from datetime import date

from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

# ===========================================================================
# Brand
# ===========================================================================
BRAND = "Lucid"
SLOGAN_EN = "Be lucid."
TAGLINE_EN = "Your verified second brain."
THESIS_EN = "AI generates infinitely. Lucid validates."

FOOTER_BRAND = f"{BRAND} · {SLOGAN_EN}"

OUT = Path(__file__).parent
DOCX = OUT / f"{BRAND}_Business_Plan.docx"
PPTX = OUT / f"{BRAND}_Pitch_Deck.pptx"

PRIMARY = RGBColor(0x1F, 0x2E, 0x5C)
ACCENT = RGBColor(0x6B, 0x8A, 0xFD)
SUCCESS = RGBColor(0x1F, 0x7A, 0x4D)
WARNING = RGBColor(0xB8, 0x6E, 0x00)
INK = RGBColor(0x1A, 0x1F, 0x2C)
GREY = RGBColor(0x6B, 0x72, 0x80)

PRIMARY_P = PRGBColor(0x1F, 0x2E, 0x5C)
ACCENT_P = PRGBColor(0x6B, 0x8A, 0xFD)
SUCCESS_P = PRGBColor(0x1F, 0x7A, 0x4D)
WARNING_P = PRGBColor(0xB8, 0x6E, 0x00)
INK_P = PRGBColor(0x1A, 0x1F, 0x2C)
GREY_P = PRGBColor(0x6B, 0x72, 0x80)
WHITE_P = PRGBColor(0xFF, 0xFF, 0xFF)
BG_P = PRGBColor(0xF6, 0xF7, 0xFB)
DARK_BG = PRGBColor(0x0A, 0x0F, 0x24)

FONT = "Calibri"


# ===========================================================================
# DOCX helpers
# ===========================================================================

def _set_run_font(run, *, font=FONT, size=11, color=INK, bold=False, italic=False):
    run.font.name = font
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.append(rFonts)
    rFonts.set(qn("w:eastAsia"), font)
    rFonts.set(qn("w:ascii"), font)
    rFonts.set(qn("w:hAnsi"), font)
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.bold = bold
    run.italic = italic


def add_para(doc, text, *, size=11, color=INK, bold=False, italic=False, align=None, space_after=4):
    p = doc.add_paragraph()
    if align is not None:
        p.alignment = align
    r = p.add_run(text)
    _set_run_font(r, size=size, color=color, bold=bold, italic=italic)
    p.paragraph_format.space_after = Pt(space_after)
    return p


def add_heading(doc, text, level=1):
    sizes = {1: 22, 2: 16, 3: 13}
    p = doc.add_paragraph()
    r = p.add_run(text)
    _set_run_font(r, size=sizes.get(level, 12), color=PRIMARY, bold=True)
    p.paragraph_format.space_before = Pt(10 if level == 1 else 6)
    p.paragraph_format.space_after = Pt(6)
    if level == 1:
        pBdr = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "8")
        bottom.set(qn("w:space"), "1")
        bottom.set(qn("w:color"), "1F2E5C")
        pBdr.append(bottom)
        p._p.get_or_add_pPr().append(pBdr)
    return p


def add_bullet(doc, text, *, size=11):
    p = doc.add_paragraph(style="List Bullet")
    p.paragraph_format.left_indent = Cm(0.6)
    p.paragraph_format.space_after = Pt(2)
    r = p.add_run(text)
    _set_run_font(r, size=size, color=INK)
    return p


def add_table(doc, rows, header=True):
    table = doc.add_table(rows=len(rows), cols=len(rows[0]))
    table.autofit = False
    table.style = "Light Grid Accent 1"
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = table.rows[ri].cells[ci]
            cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
            p = cell.paragraphs[0]
            for r in list(p.runs):
                r.text = ""
            run = p.add_run(str(cell_text))
            is_header = header and ri == 0
            _set_run_font(run, size=10, color=PRIMARY if is_header else INK, bold=is_header)
    return table


def add_page_break(doc):
    doc.add_page_break()


def add_callout(doc, text, *, size=12):
    """Left-bordered pull-quote for the single most important line of a section."""
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Cm(0.5)
    p.paragraph_format.space_before = Pt(10)
    p.paragraph_format.space_after = Pt(10)
    r = p.add_run(text)
    _set_run_font(r, size=size, color=PRIMARY, italic=True, bold=True)
    pPr = p._p.get_or_add_pPr()
    pBdr = OxmlElement("w:pBdr")
    left = OxmlElement("w:left")
    left.set(qn("w:val"), "single")
    left.set(qn("w:sz"), "18")
    left.set(qn("w:space"), "10")
    left.set(qn("w:color"), "6B8AFD")
    pBdr.append(left)
    pPr.append(pBdr)
    return p


def add_lead(doc, text):
    """Section lead paragraph - slightly larger, sets up the section."""
    return add_para(doc, text, size=11, space_after=8)


def add_body(doc, text):
    """Standard flowing body paragraph."""
    return add_para(doc, text, size=10.5, space_after=8)


# ===========================================================================
# DOCX content
# ===========================================================================

def build_docx():
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = FONT
    style.font.size = Pt(11)

    # ---- Title page ----
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(80)
    r = p.add_run(BRAND)
    _set_run_font(r, size=36, color=PRIMARY, bold=True)

    # Hero slogan
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(40)
    r = p.add_run(SLOGAN_EN)
    _set_run_font(r, size=64, color=PRIMARY, bold=True)

    # Tagline
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(50)
    r = p.add_run(TAGLINE_EN)
    _set_run_font(r, size=14, color=INK, italic=True)

    # Thesis
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(60)
    r = p.add_run(THESIS_EN)
    _set_run_font(r, size=11, color=PRIMARY, italic=True, bold=True)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Validation infrastructure for the post-AI internet.")
    _set_run_font(r, size=10, color=GREY, italic=True)

    # Document metadata
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(40)
    r = p.add_run("Business Plan - Seed Round\n")
    _set_run_font(r, size=11, color=INK, bold=True)
    r = p.add_run(f"\n{date.today().strftime('%B %Y')}\n")
    _set_run_font(r, size=10, color=GREY)
    r = p.add_run("CONFIDENTIAL - for investor review only")
    _set_run_font(r, size=9, color=GREY, italic=True)

    add_page_break(doc)

    # ===================================================================
    # 0. Executive Summary
    # ===================================================================
    add_heading(doc, "0. Executive Summary", 1)

    p = doc.add_paragraph()
    r = p.add_run("Lucid's promise is two words - ")
    _set_run_font(r, size=12, color=INK, italic=True)
    r = p.add_run(SLOGAN_EN)
    _set_run_font(r, size=14, color=PRIMARY, bold=True)
    p.paragraph_format.space_after = Pt(8)

    add_body(doc, "We live in an age of rampant misinformation, where AI generates content without limit. People consume more information than ever, remember less of it than ever, and their ability to tell what is true has, if anything, regressed. Helping users restore the lucidity of their own cognition - that is what Lucid is.")

    add_body(doc, "In one sentence, Lucid is a verified second brain. You capture information from the web and mobile the way you always have; AI organizes it into atomic facts; and you judge whether each is true in under thirty seconds. Only knowledge that passes validation remains permanently in your personal knowledge graph - and it surfaces on its own the moment you write or make a decision.")

    add_body(doc, "But Lucid's essential value is not a note-taking app. As AI prints content without limit, one resource alone grows scarce - validated truth - and Lucid is the infrastructure that accumulates it at the individual scale, and ultimately at the social scale.")

    add_callout(doc, "Validation requires human time. And time is the one resource a larger AI model cannot save for you.")

    add_body(doc, "This is where the defensive line forms. ChatGPT can absorb a memory feature in its next version, then search, then notes. But the validation labor that thousands of users accumulate, thirty seconds at a time, cannot be replicated - no matter how far model performance is pushed.")

    add_body(doc, "This venture does not pour large capital into the seed stage. Validation infrastructure is a business whose value lies in deep accumulation, not fast growth, and so we take a capital-efficient, staged approach. Phase 0 is a bootstrap of KRW 50-150M that confirms the first signal of product-market fit within six months; only when that signal appears do we advance to a Phase 1 seed round.")

    add_body(doc, "The founder built two full-stack prototypes single-handedly in six months, with no outside funding - WisdomDB, with a capture, graph, and validation workflow, and Student, a fact extraction, validation, and Q&A engine. Lucid integrates the two; the integration architecture is already designed, and a beta is possible within eight weeks.")

    add_page_break(doc)

    # ===================================================================
    # 1. Problem
    # ===================================================================
    add_heading(doc, "1. The Problem - We No Longer Remember, and Now We No Longer Trust", 1)

    add_lead(doc, "That article you found so striking last week - do you remember where you saw it? The statistic that made an argument convincing - can you recall its source? The paper a colleague recommended three months ago - what was its title?")

    add_body(doc, "For most knowledge workers, the answer to these questions is 'I don't know.' We are exposed to an average of 34 GB of information every day (UC San Diego, 2009), yet less than 5% of it survives in a form we can retrieve and reuse later. McKinsey estimates that knowledge workers spend 9.3 hours a week re-finding information they have already seen. More than a full workday a week, evaporating.")

    add_body(doc, "This is an old problem. But as of 2024, a deeper crisis has been layered on top of it.")

    add_body(doc, "The volume of AI-generated content has overtaken the volume written by humans (Europol). A substantial share of search results, social feeds, and news articles is now not written by people. Some of it is carefully engineered falsehood; far more is unsourced guesswork dressed up to look plausible. The problem is no longer finding information - it is that judging whether the information you found can be trusted has become all but impossible.")

    add_body(doc, "The two crises share a root. The human cognitive system was not designed to handle this scale, or this much uncertainty. Working memory holds only about seven items at once, and encoding into long-term memory demands conscious effort. The supply of information has grown exponentially, but the human biological capacity to validate, store, and recall it is unchanged from tens of thousands of years ago.")

    add_callout(doc, "The generation with the most access to information in human history answers the question 'what do I actually know?' with the least confidence.")

    # ===================================================================
    # 2. The Insight
    # ===================================================================
    add_heading(doc, "2. The Insight - The Real Bottleneck Is Not Storage, It Is Validation", 1)

    add_lead(doc, "Every existing answer to this problem shares the same blind spot.")

    add_body(doc, "Search engines find more information, faster. But search solves 'where is it,' not 'can I trust it.' As AI content floods in, the reliability of search results only falls.")

    add_body(doc, "Note apps - Notion, Obsidian, Evernote - let you store information. But they treat everything you save as equally true. A note you wrote three years ago based on a misunderstanding sits side by side, at equal weight, with one you carefully cross-checked yesterday. And storing and organizing is entirely manual work for the user.")

    add_body(doc, "AI assistants like ChatGPT answer your questions. But that answer is 'average knowledge' learned from the whole internet, and it is often confidently wrong. Above all, it is not your validated knowledge.")

    add_body(doc, "None of the three kinds of tool tracks the single question that matters - 'Has this information been validated, by whom, and how much?'")

    add_body(doc, "This is the core insight. The real bottleneck of the information-overload age is not storage capacity. Storage has become infinitely cheap. The real bottleneck is validation. And, decisively, validation cannot be automated. AI can generate, summarize, and rearrange information without limit, but responsibility for the final judgment - 'is this true?' - rests with a human. Asking AI to validate is like asking the author of a text whether that text is true.")

    add_callout(doc, "The next generation of knowledge tools must divide labor between AI's capacity to generate and the human capacity to judge. AI drafts; the human decides.")

    add_body(doc, "Lucid is designed precisely on that division of labor.")

    add_page_break(doc)

    # ===================================================================
    # 3. Product
    # ===================================================================
    add_heading(doc, "3. The Product - How Lucid Works", 1)

    add_lead(doc, "A concrete scene is the fastest way to explain it.")

    add_body(doc, "You are a researcher working on climate policy. Tuesday afternoon, you watch a lecture on the economics of renewable energy on YouTube. A striking claim comes up. Normally you would let it pass, or write a single line in a note app and forget it.")

    add_body(doc, "Instead, you right-click on the video and choose 'Save to Lucid.' Lucid transcribes the video's audio to text, and AI extracts three factual claims from it and shows them to you. What you did was one click; the time it took was five seconds. This is the capture stage.")

    add_body(doc, "Here is where Lucid parts ways with every other tool. You judge the three extracted claims within thirty seconds. The first has two reliable sources, so you 'accept' it. The second has weak support, so you mark it 'uncertain.' The third is merely the speaker's opinion, so you 'discard' it. This thirty seconds of judgment - Human-in-the-Loop (HITL) validation - is the heart of Lucid.")

    add_body(doc, "An accepted fact is absorbed into your personal knowledge graph. It is not merely stored - it is automatically linked to related facts already in the graph. Around entities like 'CO2,' 'electricity prices,' and 'renewable energy,' knowledge forms a web.")

    add_body(doc, "Three weeks later. You are writing a policy report. The moment you type the words 'renewable energy' into the document, the facts you saved that Tuesday surface in the Lucid sidebar. You did not open a search box. You did not strain to recall what you saved, or where. It simply surfaced - the way the human brain works. One click cites that fact, and its source follows automatically.")

    add_body(doc, "This is Lucid's four-stage loop. The friction of capture and validation never exceeds five and thirty seconds respectively, and surfacing happens before the user is even aware of it.")

    add_table(doc, [
        ["Stage", "What happens", "Who", "Friction"],
        ["Capture", "Information from web, mobile, and multimodal sources is saved in place", "User + AI", "< 5 sec"],
        ["Structure", "AI organizes it into atomic facts, entities, and relations", "AI", "Automatic"],
        ["Validate", "Truth is judged via accept / edit / discard", "User (HITL)", "< 30 sec"],
        ["Surface", "Validated facts surface automatically in the context that needs them", "AI", "User unaware"],
    ])

    # ===================================================================
    # 4. Validation Layer
    # ===================================================================
    add_heading(doc, "4. The Validation System - Four Layers of Truth", 1)

    add_lead(doc, "No fact in Lucid exists as a binary of 'stored / not stored.' Each fact carries four layers of validation marking at once, and the user sees those marks on a single screen. This is what makes Lucid validation infrastructure rather than a note app.")

    add_body(doc, "Take the fact from the earlier scenario again - 'the transition to renewable energy raises electricity prices in the short term.' Beside this fact, four validation tiers are displayed.")

    add_body(doc, "L1 . Self-validation. Did you accept this fact yourself, and how many sources did you cross-check? It is the most basic tier, and it forms the instant you capture.")

    add_body(doc, "L2 . Trust-network validation. Of the people you have chosen to trust - colleagues, researchers in your field, experts you subscribe to - how many have validated the same fact? It is shown in a form like '38 of 47 in your trust network agree (81%).'")

    add_body(doc, "L3 . System validation. What percentage of all Lucid users hold a fact consistent with this one? '73% of 6,421 users agree.' This figure is aggregated anonymously and never exposes any individual user's graph.")

    add_body(doc, "L4 . Expert validation. Has a verified expert in the relevant domain certified this fact? 'Verified by Dr. ___, PhD in energy economics.'")

    add_body(doc, "On top of this, Lucid shows a map of dissent alongside. '19% of all users (1,221 people) store the opposing fact - renewable energy lowers electricity prices in the long term. Note: this opposing fact has not passed L4 expert validation.' And provenance tracing: which report this fact came from, and which source that report came from in turn, traced back N hops.")

    add_table(doc, [
        ["Tier", "Validator", "Build priority", "Matching revenue line"],
        ["L1 Self", "The user (HITL)", "M0 - already built", "Pro subscription"],
        ["L3 System", "Anonymous aggregate of all users", "M6 - automatic once users gather", "Platform (network effect)"],
        ["L2 Trust network", "Friends, teams, subscribed users", "M12", "Team tier"],
        ["L4 Expert", "Domain expert certification", "M18+", "Marketplace / Enterprise"],
    ])

    add_callout(doc, "Other tools show you 'what you stored.' Lucid shows you 'how true it is.'")

    add_body(doc, "This transparency is the point. When users encounter a fact, they no longer vaguely wonder 'is this right?' They see, as data, what they themselves, the people they trust, the collective, and experts each say. Neither Obsidian, nor Notion, nor Mem.ai, nor ChatGPT tracks even one of these four tiers. The four tiers are also a staged market-entry path - covered in Section 9.")

    add_page_break(doc)

    # ===================================================================
    # 5. Defensibility
    # ===================================================================
    add_heading(doc, "5. Defensibility - Why This Cannot Be Copied", 1)

    add_lead(doc, "The sharpest question an investor asks an AI startup is always the same. 'What if OpenAI, or ChatGPT, simply adds this feature in its next version?'")

    add_body(doc, "For most AI-wrapper startups, this question is fatal. Their value is nothing more than a thin interface laid over a model, and if the model company chooses to, it vanishes overnight. Lucid is free of this question. The reason lies in the nature of the asset we accumulate.")

    add_body(doc, "What Lucid accumulates over time is neither model weights nor interface know-how. It is the sum of the validation judgments that thousands, tens of thousands of users have made, thirty seconds at a time. And validation requires human time. A larger model, faster inference, a longer context window can accelerate content generation - but they cannot shorten the human time it takes to judge 'is this true?'")

    add_body(doc, "So even if OpenAI adds memory and notes to ChatGPT, they do not have the graph our users have validated over years. That graph cannot be bought with money, or with model performance. It is made only with time.")

    add_body(doc, "On top of this, four moats compound. First, switching cost - for users to leave Lucid means discarding their entire cognitive asset, validated over years. Second, the compounding of personalization - Lucid's surfacing accuracy rises to fit the user's patterns of thought the more it is used, and a new entrant always starts from a blank page. Third, the scarcity of validation data - a human-validated fact graph is a kind of training data that AI companies do not have. Fourth, network effects - L3 system validation grows more precise the more users there are, and the consensus of 1,000 people and the consensus of a million are qualitatively different assets.")

    add_callout(doc, "Secure a twelve-month lead, and catching up becomes a problem not of capital but of time. And time is something no one can compress.")

    add_para(doc, "Three anticipated objections, resolved in advance.", size=11, bold=True, color=PRIMARY)

    add_body(doc, "\"Isn't this just another version of Notion or Obsidian?\" - No. Those tools are storage tools, and they assume all user input is equally true. Lucid's category is not storage but validation. We do not keep notes. We track the state of truth.")

    add_body(doc, "\"How many people will really spend 30 seconds validating?\" - Our early target is not 'everyone.' It is people who are professionally sensitive to whether information is true - researchers, analysts, journalists, policy staff. For them, thirty seconds of validation is not a cost; it is the act of turning work they already do into an asset.")

    add_body(doc, "\"Why not let AI do the validating too?\" - AI can assist validation (contradiction detection, source organization). But the moment we hand the final judgment to AI, we return to the world of 'plausibly wrong answers.' The finality of human judgment is exactly what Lucid sells.")

    add_page_break(doc)

    # ===================================================================
    # 6. Market
    # ===================================================================
    add_heading(doc, "6. The Market - Who Pays, Why, and How Much", 1)

    add_lead(doc, "Lucid's first customers are knowledge workers with a professional stake in whether information is true.")

    add_body(doc, "Concretely, four groups - academic and industry researchers, finance and strategy analysts, journalists and content creators, policy and legal practitioners. What they share is that a judgment based on wrong information comes back as a direct cost. For them, validation is not a choice but a job function. So the question 'will you spend time validating?' does not even arise for this group - they already validate, and Lucid simply turns that labor into an asset for the first time.")

    add_body(doc, "We estimate market size conservatively.")

    add_table(doc, [
        ["Segment", "Size (Year 3)", "Basis"],
        ["TAM (global knowledge workers + trust infrastructure)", "$25B", "1B knowledge workers x 1% paid conversion x $150/yr + part of the KM market"],
        ["SAM (AI-tool early adopters)", "$700M", "Korea + APAC English-speaking + US early adopters"],
        ["SOM (Year 3 target)", "$30M ARR", "4.3% of SAM"],
    ])

    add_body(doc, "Comparable valuations hint at this category's potential. The AI search company Perplexity was valued at roughly $9B in 2025, and the enterprise knowledge search company Glean at roughly $2.2B. Both companies solve the problem of 'finding information.' No company yet solves 'is this information true?' Wikipedia handles part of it through community validation, but it is non-profit, limited to English general knowledge, and not personalized.")

    add_callout(doc, "Search built a $9B company. Validation is still empty.")

    # ===================================================================
    # 7. Competition
    # ===================================================================
    add_heading(doc, "7. Competition - Obsidian, and Beyond", 1)

    add_lead(doc, "The competitor to analyze most seriously is Obsidian. Not Notion.")

    add_body(doc, "Obsidian sits at the peak of personal knowledge management (PKM) tools. With an estimated headcount of about 15, it earns an estimated $30M+ in annual revenue at an operating margin above 80%, and has never taken VC money. Every note is stored as a plain Markdown file on the user's own computer, so the data survives even if the company disappears. An ecosystem of more than 1,500 community plugins creates a powerful lock-in. After Roam Research self-destructed through pricing missteps and management turmoil, Obsidian stands at the summit of trust in the PKM market.")

    add_body(doc, "But the very structure of that strength creates Lucid's entry path. Obsidian cannot become an AI-native product, philosophically or structurally. Local-first, cloud-refusing, self-funded management - these three principles are Obsidian's identity, and at the same time they make it unable to bear server-side AI computation. In fact, Obsidian's AI features depend on external plugins, and even the database feature the company built itself took two years to ship.")

    add_body(doc, "An interesting market signal sits here. Obsidian users are assembling AI plugins themselves, and paid training programs that teach how to assemble them have even appeared. This tells us two things at once. One, demand for AI-combined knowledge management is already validated. Two, the current assembly-based solution is too cumbersome and - above all - does not address the layer of validation at all.")

    add_body(doc, "At this point Lucid's target becomes clear. We are not trying to take Obsidian's heavy users - they are people who have built workflows over years. Our market is two groups. One, the people who tried Obsidian and churned out from its complexity (an estimated several million). Two, the gap between Notion's 100M users and Obsidian's roughly 7M - the 90M-plus who wanted a PKM tool but found it too hard to even start.")

    add_table(doc, [
        ["", "General data (the internet)", "Personal data (mine)"],
        ["Validated / curated", "Wikipedia . Snopes (slow . narrow scope)", "* Lucid (the empty quadrant)"],
        ["Raw / unvalidated", "Google . Perplexity . ChatGPT", "Obsidian . Notion . Mem . Rewind"],
    ])

    add_page_break(doc)

    # ===================================================================
    # 8. Business Model
    # ===================================================================
    add_heading(doc, "8. Business Model", 1)

    add_lead(doc, "Lucid has five revenue lines, mapped to the four validation tiers.")

    add_table(doc, [
        ["Tier", "Price", "Validation tier access", "Target"],
        ["Free", "$0", "L1 + part of L3", "Acquisition . trial"],
        ["Pro", "$19/mo", "L1 + L3 + L2 preview", "Individual knowledge workers"],
        ["Team", "$39/seat", "L1 + L2 + L3", "Study groups . labs . small teams"],
        ["Enterprise", "$50-200/seat/mo", "L1-L4 + domain expert pool", "Legal . R&D . consulting . media"],
        ["Validation API", "$10K-50K/yr", "L3 + provenance API", "AI companies . search engines . publishers"],
    ])

    add_body(doc, "The last line, the Validation API, is a new business area. It lets AI companies cross-check their model's output against our validation graph; Anthropic, OpenAI, and Perplexity are potential customers.")

    add_body(doc, "The Pro tier's unit economics are healthy. LLM inference costs about $3 per user per month (on Sonnet, lower with prompt caching), and at a $19 subscription the gross margin is about 78%. We estimate 12-month retention at 65% - because the lock-in of validation assets exceeds the PKM category average. On that basis LTV is about $310, target CAC is $90, and so LTV/CAC is 3.4x, above the SV seed-stage baseline of 3x. CAC payback is about five months.")

    add_body(doc, "From Year 2, L3 network effects naturally drive CAC down - as users grow, validation-consensus data grows richer, and that raises product value, which grows referral-driven acquisition.")

    # ===================================================================
    # 9. Execution & Funding
    # ===================================================================
    add_heading(doc, "9. Execution - A Staged Approach and Funding Strategy", 1)

    add_lead(doc, "This venture starts small and expands in stages. Validation infrastructure is a business whose value lies in deep accumulation, not fast growth, and since the accumulation of user validation labor is itself the value, the size of the raise is not the decisive variable. A standard $1.5M SV seed round is not appropriate for this stage.")

    add_para(doc, "Phase 0 - Bootstrap + Angel (0-6 months, current round)", size=12, bold=True, color=PRIMARY)
    add_body(doc, "The target is KRW 50-150M - founder capital plus two or three angels. Operating solo with one part-time engineer added, we confirm the first signal of product-market fit within six months. Only when we hit 50 paying users, retention above 60%, and an NPS above 40 at month six do we enter Phase 1. If we fall short, we re-examine the hypothesis. Below is the six-month cost breakdown.")

    add_table(doc, [
        ["Item", "6-month total (USD)", "Note"],
        ["Anthropic API", "$1,800", "Haiku-centric, prompt caching"],
        ["Infrastructure (vector DB . hosting . Neo4j)", "$930", "Aggressive use of free tiers early on"],
        ["Domain + transactional email", "$120", ""],
        ["Design (one-off outsourcing)", "$5,000", "Logo . landing . UI setup"],
        ["Legal (one-off)", "$3,000", "Korean attorney counsel (terms . privacy)"],
        ["Marketing experiments", "$1,800", "A/B testing . content"],
        ["Part-time engineer (20h/week)", "$24,000", "Infrastructure . integration work"],
        ["Total", "~$36,650 (about KRW 50M)", ""],
    ])

    add_para(doc, "Phase 1 - Korean Seed (6-18 months)", size=12, bold=True, color=PRIMARY)
    add_body(doc, "Once the Phase 0 signal is confirmed, we run a Korean seed round of KRW 300M-1B - Sparklabs, Primer, Mashup Angels, and Strong Ventures are the targets. At this stage the founder converts to full-time and hires a founding engineer and a designer. We move GTM experiments into full gear, enter the Validation API beta, and begin cold outreach to potential partners such as Anthropic. The month-18 targets are $30K in monthly revenue and three enterprise LOIs.")

    add_para(doc, "Phase 2 - Series A or Profitability (18 months+)", size=12, bold=True, color=PRIMARY)
    add_body(doc, "At this point we keep optionality. If we reach $1M ARR and net revenue retention above 110%, we move toward a Series A; even if not, once we reach $2M ARR at profitability we can choose to run independently, as Obsidian does. Both paths favor the founder, and neither is forced.")

    add_table(doc, [
        ["", "M6", "M12", "M18", "M24", "M36"],
        ["Free users", "500", "3,500", "15,000", "50,000", "200,000"],
        ["Paying users", "50", "400", "1,500", "5,000", "18,000"],
        ["Enterprise / API", "0", "0", "3", "20", "80"],
        ["MRR (USD)", "$1K", "$8K", "$30K", "$100K", "$360K"],
        ["ARR (USD)", "$12K", "$96K", "$360K", "$1.2M", "$4.3M"],
        ["Gross margin", "60%", "70%", "75%", "78%", "80%"],
    ])

    add_page_break(doc)

    # ===================================================================
    # 10. Founder
    # ===================================================================
    add_heading(doc, "10. The Founder", 1)

    add_lead(doc, "The founder begins from Carnegie Mellon University's AI lineage. SCHOLAR (Carbonell, 1970) and the Student proposal (2019) - both intellectual roots of Lucid - came out of this university.")

    add_body(doc, "But more important than credentials is the evidence of execution. The founder shipped two complete full-stack products single-handedly, in six months, with no outside funding. WisdomDB has a browser extension and a mobile PWA, a multimodal extractor, a Neo4j graph, and a HITL workflow. Student has a fact extraction, validation, and Q&A engine. Lucid integrates the mechanisms of the two, and the integration design is already done.")

    add_callout(doc, "Shipped twice with no funding. The next question is not 'can it be built' but 'how far will it go.'")

    add_body(doc, "Investment is not used to make 0 into 1. It is used to grow an already-proven 1 into 10. The hiring plan is staged - in Phase 1 we hire a founding engineer (with vector and graph search experience) and a product designer, and a GTM lead later in Phase 1. For advisors we target a CMU LTI professor, someone from Wikipedia or Glean, and a Korean EdTech serial founder.")

    # ===================================================================
    # 11. Risks
    # ===================================================================
    add_heading(doc, "11. Risks and Mitigation", 1)

    add_lead(doc, "The principal risks and their mitigation strategies are as follows.")

    add_table(doc, [
        ["Risk", "Severity", "Mitigation"],
        ["Retrieval latency", "High", "Vector + graph hybrid index, edge caching, target under 200ms"],
        ["Privacy concerns", "High", "L1 local-mode option, L3 anonymous aggregate only, BYOK support"],
        ["Obsidian + AI plugins maturing", "Medium", "Assembly-based solutions cannot implement L2-L4 validation tiers - our native advantage"],
        ["Perplexity / OpenAI entering validation", "Medium", "They are web-first by DNA - personal graphs and HITL are not their ICP"],
        ["Failure to form user habits", "Medium", "Early ICP limited to people who already validate as a job function - minimal friction"],
        ["LLM price volatility", "Low", "Multi-vendor backup + caching + BYOK"],
    ])

    # ===================================================================
    # 12. Vision
    # ===================================================================
    add_heading(doc, "12. Vision - The Age of Validated Knowledge", 1)

    add_body(doc, "In 1945, Vannevar Bush proposed the 'Memex' - a machine for an individual to link, store, and instantly retrieve all of their knowledge. In 1968, Doug Engelbart demonstrated part of that vision. In 2001, Wikipedia institutionalized 'collectively validated knowledge' for the first time.")

    add_body(doc, "Across this 80-year lineage there has always been the same empty box - a system to accumulate, frictionlessly, the knowledge an individual has validated, and to retrieve it when needed. The Memex did not address validation; Wikipedia did not address the individual. Lucid is that empty box.")

    add_body(doc, "And that empty box can, now, be filled for the first time. LLMs have removed the friction of capture and structuring; graph and vector databases have made instant surfacing possible; and above all - the flood of AI content has, at last, made validated truth a scarce and valuable resource.")

    add_body(doc, "The end picture we envision is this. A generation from now, people will have their own validated knowledge graph the way they have an email account or a calendar. Faced with a claim, they will know at once what in their graph it agrees with and what it conflicts with. Individuals' graphs, merged anonymously, let humanity share a map of validated truth - not one shaped by ads or algorithms.")

    add_callout(doc, "That is the world Lucid sets out to build. And the first words of that world are two - Be lucid.")

    # ===================================================================
    # Appendix A
    # ===================================================================
    add_heading(doc, "Appendix A - Intellectual Lineage", 1)
    add_body(doc, "Lucid sits at the intersection of two long traditions. One is the 70-year lineage of Intelligence Augmentation; the other is the lineage of attempts to institutionalize validation.")

    add_para(doc, "The Intelligence Augmentation lineage", size=11, bold=True, color=PRIMARY)
    add_table(doc, [
        ["Year", "Contribution", "Figure"],
        ["1945", "As We May Think (Memex)", "Vannevar Bush"],
        ["1960", "Man-Computer Symbiosis", "J.C.R. Licklider"],
        ["1962", "Augmenting Human Intellect", "Doug Engelbart"],
        ["1965", "Hypertext / Xanadu", "Ted Nelson"],
        ["1970", "SCHOLAR (conversational CAI)", "Jaime R. Carbonell Sr."],
        ["2019", "The Student proposal", "CMU (Carbonell school)"],
    ])

    add_para(doc, "The Institutional Validation lineage", size=11, bold=True, color=PRIMARY)
    add_table(doc, [
        ["Year", "Contribution", "Validation mechanism"],
        ["~1660", "Royal Society peer review", "Authority + reproduction"],
        ["2001", "Wikipedia", "Community consensus + edit history"],
        ["2008", "Stack Overflow", "Collective voting + reputation"],
        ["2021", "Community Notes (Birdwatch)", "Algorithm-mediated collective consensus"],
        ["2026", "Lucid", "AI-native + 4-tier (self / trust network / system / expert) + provenance tracing"],
    ])

    # ===================================================================
    # Appendix B
    # ===================================================================
    add_heading(doc, "Appendix B - References", 1)
    add_bullet(doc, "Bush, V. (1945). As We May Think. The Atlantic Monthly.")
    add_bullet(doc, "Licklider, J. C. R. (1960). Man-Computer Symbiosis. IRE Trans. Human Factors.")
    add_bullet(doc, "Engelbart, D. C. (1962). Augmenting Human Intellect. SRI Report.")
    add_bullet(doc, "Miller, G. A. (1956). The Magical Number Seven, Plus or Minus Two.")
    add_bullet(doc, "Collins, A. M. & Loftus, E. F. (1975). A spreading-activation theory of semantic processing.")
    add_bullet(doc, "Baddeley, A. D. & Hitch, G. (1974). Working memory.")
    add_bullet(doc, "Tulving, E. & Thomson, D. (1973). Encoding specificity.")
    add_bullet(doc, "Johnson, M. K. et al. (1993). Source monitoring.")
    add_bullet(doc, "Sperber, D. et al. (2010). Epistemic vigilance. Mind & Language.")
    add_bullet(doc, "Carbonell, J. R. (1970). AI in CAI. IEEE Trans. Man-Machine Systems.")
    add_bullet(doc, "McKinsey Global Institute (2023). Productivity reports.")
    add_bullet(doc, "Europol (2022). Facing Reality? Law Enforcement and the Challenge of Deepfakes.")
    add_bullet(doc, "UC San Diego (2009). How Much Information.")

    # ---- Closing slogan ----
    add_page_break(doc)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(180)
    r = p.add_run(SLOGAN_EN)
    _set_run_font(r, size=56, color=PRIMARY, bold=True)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(60)
    r = p.add_run(f"{BRAND}  .  {TAGLINE_EN}")
    _set_run_font(r, size=11, color=GREY)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("END OF DOCUMENT")
    _set_run_font(r, size=9, color=GREY, italic=True)

    return _safe_save(doc, DOCX)


# ===========================================================================
# PPTX helpers
# ===========================================================================

def _set_textbox(tf, lines, *, default_size=18, default_color=INK_P, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = PPt(opts.get("size", default_size))
        r.font.bold = opts.get("bold", False)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", default_color)


def add_title_slide(prs, *, brand, slogan_en, tagline_en, thesis_en, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    # Brand mark (top, smaller)
    tb = slide.shapes.add_textbox(PInches(0.8), PInches(0.6), PInches(11.5), PInches(0.6))
    _set_textbox(tb.text_frame, [(brand, {"size": 32, "bold": True, "color": PRGBColor(0xBF, 0xC9, 0xE8)})])

    # HERO slogan (center, biggest)
    tb = slide.shapes.add_textbox(PInches(0.8), PInches(2.4), PInches(11.5), PInches(2.0))
    _set_textbox(tb.text_frame, [
        (slogan_en, {"size": 96, "bold": True, "color": WHITE_P, "align": PP_ALIGN.CENTER}),
    ])

    # Tagline
    tb = slide.shapes.add_textbox(PInches(0.8), PInches(4.4), PInches(11.5), PInches(0.5))
    _set_textbox(tb.text_frame, [
        (tagline_en, {"size": 20, "color": PRGBColor(0x9F, 0xAB, 0xCE), "italic": True, "align": PP_ALIGN.CENTER}),
    ])

    # Thesis (smaller, accent)
    tb = slide.shapes.add_textbox(PInches(0.8), PInches(5.3), PInches(11.5), PInches(0.4))
    _set_textbox(tb.text_frame, [
        (thesis_en, {"size": 14, "color": ACCENT_P, "italic": True, "align": PP_ALIGN.CENTER}),
    ])

    # Footer
    tb = slide.shapes.add_textbox(PInches(0.8), PInches(6.95), PInches(11.5), PInches(0.3))
    _set_textbox(tb.text_frame, [
        (subtitle, {"size": 11, "color": PRGBColor(0x6B, 0x77, 0x9C), "align": PP_ALIGN.CENTER}),
    ])
    return slide


def add_content_slide(prs, *, header, title_en, body_blocks=None, footer=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE_P
    bg.line.fill.background()
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, PInches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_P
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(PInches(0.6), PInches(0.25), PInches(12), PInches(0.4))
    _set_textbox(tb.text_frame, [(header.upper(), {"size": 11, "bold": True, "color": ACCENT_P})])

    tb = slide.shapes.add_textbox(PInches(0.6), PInches(0.65), PInches(12), PInches(0.9))
    _set_textbox(tb.text_frame, [(title_en, {"size": 32, "bold": True, "color": INK_P})])

    if body_blocks:
        tb = slide.shapes.add_textbox(PInches(0.7), PInches(1.75), PInches(12), PInches(4.9))
        _set_textbox(tb.text_frame, body_blocks, default_size=16, default_color=INK_P)

    if footer:
        tb = slide.shapes.add_textbox(PInches(0.6), PInches(6.85), PInches(12), PInches(0.3))
        _set_textbox(tb.text_frame, [(footer, {"size": 10, "color": GREY_P})])
    return slide


def add_table_slide(prs, *, header, title_en, table_data=None, col_widths=None, footer=None):
    slide = add_content_slide(prs, header=header, title_en=title_en, body_blocks=None, footer=footer)
    rows = len(table_data)
    cols = len(table_data[0])
    table = slide.shapes.add_table(rows, cols, PInches(0.7), PInches(1.9), PInches(12), PInches(0.55 * rows)).table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = PInches(w)
    for ri, row in enumerate(table_data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = PPt(13)
                    if ri == 0:
                        r.font.bold = True
                        r.font.color.rgb = WHITE_P
                    else:
                        r.font.color.rgb = INK_P
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY_P
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE_P if ri % 2 == 1 else BG_P
    return slide


# ===========================================================================
# PPTX content
# ===========================================================================

def build_pptx():
    prs = Presentation()
    prs.slide_width = PInches(13.33)
    prs.slide_height = PInches(7.5)

    today = date.today().strftime("%B %Y")

    # 1. Title
    add_title_slide(
        prs,
        brand=BRAND,
        slogan_en=SLOGAN_EN,
        tagline_en=TAGLINE_EN,
        thesis_en=THESIS_EN,
        subtitle=f"Seed Round  .  {today}  .  CONFIDENTIAL",
    )

    # 2. Problem
    add_content_slide(
        prs,
        header="01 . Problem",
        title_en="The crisis isn't volume. It's validation.",
        body_blocks=[
            ("·  34 GB of information a day - under 5% is retained", {"size": 18}),
            ("·  9.3 hours a week re-finding information already seen (McKinsey)", {"size": 18}),
            ("·  In 2024 AI content overtook human content (Europol) - telling true from false is near-impossible", {"size": 18}),
            ("·  Even captured into Notion or Pocket, the revisit rate is under 8%", {"size": 18}),
            ("·  ChatGPT answers with the internet's general knowledge - not validated fact", {"size": 18}),
            ("", {"size": 12}),
            ("AI accelerates generation. It cannot accelerate human validation.", {"size": 18, "color": PRIMARY_P, "bold": True, "italic": True}),
            ("The gap between them is our market.", {"size": 18, "color": PRIMARY_P, "bold": True, "italic": True}),
        ],
        footer=f"{BRAND} . {SLOGAN_EN} . 2 / 16",
    )

    # 3. Why Now
    add_table_slide(
        prs,
        header="02 . Why Now",
        title_en="Validation becomes the scarce resource.",
        table_data=[
            ["Force", "What changed", "Implication"],
            ["LLM cost", "Haiku 4.5 - ~$0.50 / user / month", "A personal cognition layer is viable"],
            ["AI content flood", "2024: AI content > human content", "A validation market forms"],
            ["Vector + graph DB maturity", "Qdrant . Neo4j stable", "Hybrid retrieval is possible"],
            ["Multimodal models", "Vision captures video and images", "Capture friction drops to zero"],
        ],
        col_widths=[2.5, 5.5, 4.0],
        footer=f"{BRAND} . {SLOGAN_EN} . 3 / 16",
    )

    # 4. The Insight
    add_content_slide(
        prs,
        header="03 . The Insight",
        title_en="AI generates infinitely. Lucid validates.",
        body_blocks=[
            ("", {"size": 10}),
            ("The tragedy of other AI startups: when the model improves, the wrapper disappears.", {"size": 18, "color": GREY_P}),
            ("What Lucid accumulates: human validation time. The one resource AI cannot save for you.", {"size": 18, "color": PRIMARY_P, "bold": True}),
            ("", {"size": 16}),
            ("Two lenses:", {"size": 17, "bold": True, "color": PRIMARY_P}),
            ("·  Consumer  ->  \"Your verified second brain\"", {"size": 17}),
            ("·  Investor  ->  \"Validation infrastructure for the post-AI internet\"", {"size": 17}),
            ("", {"size": 12}),
            ("What Wikipedia took 20 years to build - with AI tools, at 30 seconds of friction per user, down to the individual.", {"size": 15, "italic": True, "color": ACCENT_P}),
        ],
        footer=f"{BRAND} . {SLOGAN_EN} . 4 / 16",
    )

    # 5. Demo
    add_content_slide(
        prs,
        header="04 . Demo",
        title_en="60 seconds. Capture -> 4-tier validate -> Surface.",
        body_blocks=[
            ("0:00   Watching a YouTube lecture", {"size": 17}),
            ("0:10   Right-click -> Lucid -> AI runs speech-to-text -> extracts 3 facts", {"size": 17}),
            ("0:20   L1 self-validation (5-sec HITL) - accept 2, tag 1 \"uncertain\"", {"size": 17}),
            ("0:25   L3 system: 4,200 others stored the same facts (89% consensus) - shown automatically", {"size": 17, "color": ACCENT_P}),
            ("0:35   Start writing -> validated facts surface in the sidebar (with L1+L3 check badges)", {"size": 17, "bold": True, "color": PRIMARY_P}),
            ("0:50   One click cites it - source and validation status attached automatically", {"size": 17}),
            ("Later   On a contradictory claim on Twitter -> alert comparing it against your graph and L3 consensus", {"size": 17, "color": WARNING_P}),
            ("", {"size": 12}),
            ("[ 60-second demo video placeholder - Loom link ]", {"size": 13, "color": GREY_P, "italic": True}),
        ],
        footer=f"{BRAND} . {SLOGAN_EN} . 5 / 16",
    )

    # 6. THE VALIDATION LAYER (centerpiece)
    add_content_slide(
        prs,
        header="05 . The Validation Layer",
        title_en="Four tiers. One transparent badge.",
        body_blocks=[
            ("Fact: \"The transition to renewable energy raises electricity prices in the short term.\"", {"size": 15, "bold": True}),
            ("", {"size": 6}),
            ("L1  Self-validated        You (2 sources cross-checked)", {"size": 15, "color": SUCCESS_P}),
            ("L2  Peer-validated        38 of 47 in your trust network agree (81%)", {"size": 15, "color": SUCCESS_P}),
            ("L3  System consensus     73% of 6,421 users agree", {"size": 15, "color": SUCCESS_P}),
            ("L4  Expert-verified       Verified by Dr. Park (PhD, energy economics)", {"size": 15, "color": SUCCESS_P}),
            ("", {"size": 4}),
            ("Contradiction map", {"size": 15, "color": WARNING_P, "bold": True}),
            ("    · 19% of users (n=1,221) store the opposing fact: \"lowers prices in the long term\"", {"size": 13, "color": WARNING_P}),
            ("    · the opposing fact is not L4-verified", {"size": 13, "color": WARNING_P}),
            ("", {"size": 4}),
            ("Provenance trail (4 hops):  BloombergNEF -> IEA -> government report -> original paper", {"size": 13, "color": GREY_P}),
            ("", {"size": 6}),
            ("This transparency exists in no other tool - Obsidian, Notion, Mem, ChatGPT. It defines our category.", {"size": 14, "bold": True, "color": PRIMARY_P, "italic": True}),
        ],
        footer=f"{BRAND} . {SLOGAN_EN} . 6 / 16  .  Centerpiece",
    )

    # 7. Mechanism
    add_table_slide(
        prs,
        header="06 . Mechanism",
        title_en="Cognitive science + epistemic vigilance.",
        table_data=[
            ["Cognitive process", "Theory", "Lucid implementation"],
            ["Spreading activation", "Collins & Loftus (1975)", "Graph traversal + vector ANN"],
            ["Working <-> long-term memory", "Baddeley & Hitch (1974)", "Context-aware inline surfacing"],
            ["Encoding specificity", "Tulving (1973)", "Capture context preserved"],
            ["Source monitoring", "Johnson et al. (1993)", "Permanent provenance edges"],
            ["Epistemic vigilance", "Sperber et al. (2010)", "4-tier validation badge"],
        ],
        col_widths=[3.2, 3.3, 5.5],
        footer=f"{BRAND} . {SLOGAN_EN} . 7 / 16",
    )

    # 8. Killer Use Cases
    add_content_slide(
        prs,
        header="07 . Killer Use Cases",
        title_en="Search becomes surfacing. Notes become evidence.",
        body_blocks=[
            ("·  Drafting Assistant - validated facts surface as you write; L1-L4 badges shown on citation", {"size": 16}),
            ("·  Truth Lens - meet a claim while browsing, see its L3 consensus against your graph instantly", {"size": 16}),
            ("·  Decision Support - when deciding, see related nodes, past judgments, and 4-tier status at a glance", {"size": 16}),
            ("·  Re-discovery - natural-language search + validation-tier filter (\"show me only L3 80%+\")", {"size": 16}),
            ("·  Conflict Resolver - when a new capture conflicts with an existing fact, auto-alert + both sides' evidence", {"size": 16}),
        ],
        footer=f"{BRAND} . {SLOGAN_EN} . 8 / 16",
    )

    # 9. Market
    add_table_slide(
        prs,
        header="08 . Market",
        title_en="Trust & truth infrastructure - a newly emerging category.",
        table_data=[
            ["Comparable", "Valuation / ARR", "What they validate"],
            ["Perplexity", "$9B (2025)", "No validation - AI search"],
            ["Glean", "$2.2B", "No validation - enterprise search"],
            ["Wikipedia (non-profit)", "270M monthly users", "Community-validated, English general knowledge"],
            ["Snopes / FactCheck", "~$10M ARR", "Expert-validated, narrow scope"],
            ["Community Notes (X)", "Built into X", "Crowd-consensus algorithm"],
            ["TAM / SAM / SOM (3yr)", "$25B / $700M / $30M ARR", "1B knowledge workers x 1% x $150"],
        ],
        col_widths=[3.0, 2.7, 6.6],
        footer=f"{BRAND} . {SLOGAN_EN} . 9 / 16",
    )

    # 10. Competition - Obsidian-centric
    add_content_slide(
        prs,
        header="09 . Competition",
        title_en="Obsidian is the strongest comparator. Here's the gap.",
        body_blocks=[
            ("Obsidian's strengths:", {"size": 14, "bold": True, "color": PRIMARY_P}),
            ("·  Local-first DNA . trust capital . 1,500+ plugins . ~15 staff, est. $30M ARR, 80% margin", {"size": 13}),
            ("", {"size": 6}),
            ("Obsidian's limits = our entry path:", {"size": 14, "bold": True, "color": PRIMARY_P}),
            ("·  30%+ churn in week one of the learning curve . AI-native impossible (philosophically & structurally) . no validation mechanism . no collaboration", {"size": 13}),
            ("", {"size": 6}),
            ("Market signal:", {"size": 14, "bold": True, "color": PRIMARY_P}),
            ("·  Obsidian + AI-plugin assembly workflows + paid training programs = AI-PKM demand is validated", {"size": 13}),
            ("·  but assembly-based solutions cannot implement a validation layer - our native advantage", {"size": 13}),
            ("", {"size": 8}),
            ("Our ICP (not Obsidian heavy users):", {"size": 14, "bold": True, "color": PRIMARY_P}),
            ("·  Obsidian try-and-churn users (4.5M) + knowledge workers who never even tried PKM (93M)", {"size": 13}),
            ("·  those aware of the AI trust crisis (researchers . journalists . legal . R&D)", {"size": 13}),
            ("", {"size": 4}),
            ("                  General data            Personal data", {"size": 12, "color": GREY_P}),
            ("Validated     Wikipedia . Snopes      * LUCID (empty)", {"size": 14, "bold": True}),
            ("Unfiltered     Google . ChatGPT       Obsidian . Notion . Mem . Rewind", {"size": 12, "color": GREY_P}),
        ],
        footer=f"{BRAND} . {SLOGAN_EN} . 10 / 16",
    )

    # 11. Wedge & Expansion (4 layers as ladder)
    add_table_slide(
        prs,
        header="10 . Wedge & Expansion",
        title_en="Four validation layers = four expansion stages.",
        table_data=[
            ["Stage", "Layer", "Tier", "When"],
            ["1. Personal cognition aid", "L1 Self", "Pro $19/mo", "2026"],
            ["2. System consensus", "L3 System (network effect)", "Platform layer", "2026-27"],
            ["3. Trust-network validation", "L2 Peer", "Team $39/seat", "2027"],
            ["4. Expert validation", "L4 Expert", "Marketplace / Enterprise", "2028"],
        ],
        col_widths=[3.4, 3.6, 3.3, 1.7],
        footer=f"{BRAND} . {SLOGAN_EN} . 11 / 16  .  L3 is decisive - it emerges automatically once users gather, and cannot be copied",
    )

    # 12. Business Model
    add_table_slide(
        prs,
        header="11 . Business Model",
        title_en="4 validation layers = 5 revenue lines.",
        table_data=[
            ["Tier", "Price", "Validation access", "Target"],
            ["Free", "$0", "L1 + part of L3", "Acquisition"],
            ["Pro", "$19/mo", "L1 + L3 + L2 preview", "Individual knowledge workers"],
            ["Team", "$39/seat", "L1 + L2 + L3", "Study groups . labs . teams"],
            ["Enterprise", "$50-200/seat/mo", "L1-L4 + domain experts", "Legal . R&D . consulting"],
            ["Validation API", "$10K-50K/yr", "L3 + provenance API", "AI companies . search . publishers"],
        ],
        col_widths=[2.0, 2.3, 3.8, 3.9],
        footer=f"{BRAND} . {SLOGAN_EN} . 12 / 16  .  LTV $310 . CAC $90 . LTV/CAC 3.4x",
    )

    # 13. Defensibility
    add_content_slide(
        prs,
        header="12 . Defensibility",
        title_en="Five moats compound.",
        body_blocks=[
            ("1.   Switching cost - leave, and your validated cognitive asset + trust network + provenance trail are all gone", {"size": 15}),
            ("2.   Compound personalization - the ranker fits you the more you use it. New entrants face cold-start", {"size": 15}),
            ("3.   HITL-validated data - a labeled asset Anthropic / OpenAI do not have. Partnership / data-licensing optionality", {"size": 15}),
            ("4.   Provenance graph - every fact permanently linked to its source. A verifiability LLMs cannot imitate", {"size": 15}),
            ("5.   Network effect (L3) - more users, richer L3 system consensus. Latecomers cannot catch up", {"size": 15, "color": ACCENT_P, "bold": True}),
            ("", {"size": 10}),
            ("Not an \"AI wrapper.\" We lock the user's *validation labor* into a workflow LLMs cannot reach.", {"size": 15, "italic": True, "color": PRIMARY_P}),
            ("With a 12-month lead, the structure becomes impossible to catch.", {"size": 15, "italic": True, "color": PRIMARY_P}),
        ],
        footer=f"{BRAND} . {SLOGAN_EN} . 13 / 16",
    )

    # 14. Traction
    add_content_slide(
        prs,
        header="13 . Traction",
        title_en="Shipped twice before raising.",
        body_blocks=[
            ("Capture + graph + HITL layer (WisdomDB)", {"size": 16, "bold": True}),
            ("        Browser extension + mobile PWA + YouTube / PDF / image / audio extractors", {"size": 13, "color": GREY_P}),
            ("        Neo4j ontology + accept / edit / discard workflow", {"size": 13, "color": GREY_P}),
            ("", {"size": 6}),
            ("Fact extraction . validation . answering engine (Student)", {"size": 16, "bold": True}),
            ("        The full learning-AI loop of the CMU SCHOLAR lineage (Carbonell, 1970)", {"size": 13, "color": GREY_P}),
            ("        The core component of a \"personal LLM brain\" that reasons using validated facts only", {"size": 13, "color": GREY_P}),
            ("", {"size": 6}),
            (f"->   Integrating the two systems = {BRAND}'s cognition aid + 4-tier validation loop. Integrated beta within 8 weeks.", {"size": 16, "bold": True, "color": PRIMARY_P}),
            ("", {"size": 8}),
            ("[ N beta signups - update once recruitment data is in ]", {"size": 12, "color": GREY_P, "italic": True}),
        ],
        footer=f"{BRAND} . {SLOGAN_EN} . 14 / 16",
    )

    # 15. The Ask - staged funding
    add_table_slide(
        prs,
        header="14 . The Ask",
        title_en="Staged funding. Phase 0 first.",
        table_data=[
            ["Phase", "Timeline", "Raise", "Use of funds", "Gate to advance"],
            ["Phase 0 - Bootstrap + Angel", "0-6 mo (now)", "~$37K (KRW 50-150M)", "Solo + 1 part-time engineer; first PMF signal", "50 paying users . 60%+ retention . NPS 40+"],
            ["Phase 1 - Korean Seed", "6-18 mo", "KRW 300M-1B", "Founding engineer + designer; GTM + Validation API beta", "$30K MRR . 3 enterprise LOIs"],
            ["Phase 2 - Series A / Profit", "18 mo+", "Optional", "Scale the team; keep optionality", "$1M ARR + 110% NRR  -or-  $2M ARR profitable"],
        ],
        col_widths=[2.7, 1.5, 2.0, 3.4, 2.4],
        footer=f"{BRAND} . {SLOGAN_EN} . 15 / 16  .  We advance only when the gate is met - capital follows proof, not the reverse",
    )

    # 16. Vision (closing with hero slogan)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    tb = slide.shapes.add_textbox(PInches(0.8), PInches(1.0), PInches(11.5), PInches(2.5))
    _set_textbox(tb.text_frame, [
        ("Bush proposed it in 1945.", {"size": 24, "color": PRGBColor(0xBF, 0xC9, 0xE8)}),
        ("Engelbart demoed it in 1968.", {"size": 24, "color": PRGBColor(0xBF, 0xC9, 0xE8)}),
        ("Wikipedia institutionalized it in 2001.", {"size": 24, "color": PRGBColor(0xBF, 0xC9, 0xE8)}),
        ("Lucid ships it in 2026.", {"size": 30, "color": WHITE_P, "bold": True}),
    ])

    # Big closing slogan
    tb = slide.shapes.add_textbox(PInches(0.8), PInches(4.3), PInches(11.5), PInches(1.5))
    _set_textbox(tb.text_frame, [
        (SLOGAN_EN, {"size": 96, "bold": True, "color": WHITE_P, "align": PP_ALIGN.CENTER}),
    ])

    tb = slide.shapes.add_textbox(PInches(0.8), PInches(6.0), PInches(11.5), PInches(0.5))
    _set_textbox(tb.text_frame, [
        (TAGLINE_EN, {"size": 20, "color": PRGBColor(0xBF, 0xC9, 0xE8), "italic": True, "align": PP_ALIGN.CENTER}),
    ])

    tb = slide.shapes.add_textbox(PInches(0.8), PInches(6.95), PInches(11.5), PInches(0.3))
    _set_textbox(tb.text_frame, [
        (f"{BRAND}  .  Validation infrastructure for the post-AI internet.   .   16 / 16",
         {"size": 11, "color": PRGBColor(0x6B, 0x77, 0x9C), "align": PP_ALIGN.CENTER}),
    ])

    return _safe_save(prs, PPTX)


# ===========================================================================
# Run
# ===========================================================================

def _safe_save(obj, path):
    """Save to path; if the file is locked (open in Office), save a _NEW variant."""
    try:
        obj.save(str(path))
        return path
    except PermissionError:
        alt = path.with_name(path.stem + "_NEW" + path.suffix)
        obj.save(str(alt))
        print(f"  ({path.name} is open - saved as {alt.name} instead)")
        return alt


if __name__ == "__main__":
    # Clean up old branded files (skip silently if locked by an open Word/PPT)
    for prefix in ("Praxis_", "Augment_"):
        for stale in list(OUT.glob(f"{prefix}*.docx")) + list(OUT.glob(f"{prefix}*.pptx")):
            try:
                stale.unlink()
            except PermissionError:
                print(f"  (skip locked file: {stale.name} - close Word/PPT to clean)")

    d = build_docx()
    p = build_pptx()
    print(f"Built {d.name}  ({d.stat().st_size:,} bytes)")
    print(f"Built {p.name}  ({p.stat().st_size:,} bytes)")
